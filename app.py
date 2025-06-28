import streamlit as st
import json
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import base64
from datetime import datetime
import re

# Page configuration
st.set_page_config(
    page_title="AI PPT Generator ",
   # page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)



# Initialize session state
if 'generated_content' not in st.session_state:
    st.session_state.generated_content = None
if 'pptx_file' not in st.session_state:
    st.session_state.pptx_file = None

# Header
st.markdown("""
<div class="main-header">
    <h1> AI PowerPoint Generator</h1>
    <p><span class="gemini-badge"> Gemini 2.0 Flash</span></p>
    <p>Generate professional presentations with Google's latest AI</p>
</div>
""", unsafe_allow_html=True)

# Sidebar for API configuration
st.sidebar.header(" Configuration")
st.sidebar.markdown("### Google Gemini API")
api_key = st.sidebar.text_input(
    "Gemini API Key", 
    type="password",
    help="Get your API key from Google AI Studio: https://makersuite.google.com/app/apikey"
)


# Main content area
col1, col2 = st.columns([2, 1])

with col1:
    st.header(" Presentation Details")
    
    # Topic input
    topic = st.text_input(
        "Presentation Topic",
        placeholder="e.g., Artificial Intelligence in Healthcare",
        help="Enter the main topic for your presentation"
    )
    
    # Additional details
    description = st.text_area(
        "Additional Details (Optional)",
        placeholder="Provide any specific points, audience information, or requirements...",
        height=100
    )
    
    # Presentation settings
    st.subheader(" Presentation Settings")
    
    col_settings1, col_settings2 = st.columns(2)
    
    with col_settings1:
        num_slides = st.slider("Number of Slides", min_value=3, max_value=15, value=7)
        presentation_style = st.selectbox(
            "Presentation Style",
            ["Professional", "Creative", "Academic", "Casual", "Corporate", "Technical"]
        )
    
    with col_settings2:
        audience = st.selectbox(
            "Target Audience",
            ["General", "Students", "Executives", "Technical", "Marketing", "Healthcare", "Education"]
        )
        include_images = st.checkbox("Include Image Suggestions", value=True)

with col2:
    st.header(" Design Options")
    
    color_scheme = st.selectbox(
        "Color Scheme",
        ["Google Blue", "Professional Blue", "Green Nature", "Red Dynamic", "Purple Creative", "Orange Warm", "Monochrome"]
    )
    
    font_size = st.selectbox(
        "Font Size",
        ["Small (18pt)", "Medium (24pt)", "Large (32pt)"],
        index=1
    )
    
    template_style = st.selectbox(
        "Template Style",
        ["Modern", "Classic", "Minimalist", "Bold", "Corporate"]
    )

# AI Content Generation Functions
def get_color_scheme(scheme_name):
    schemes = {
        "Google Blue": {"primary": RGBColor(66, 133, 244), "secondary": RGBColor(52, 168, 83)},
        "Professional Blue": {"primary": RGBColor(52, 152, 219), "secondary": RGBColor(41, 128, 185)},
        "Green Nature": {"primary": RGBColor(52, 168, 83), "secondary": RGBColor(46, 204, 113)},
        "Red Dynamic": {"primary": RGBColor(234, 67, 53), "secondary": RGBColor(219, 68, 55)},
        "Purple Creative": {"primary": RGBColor(155, 89, 182), "secondary": RGBColor(142, 68, 173)},
        "Orange Warm": {"primary": RGBColor(251, 188, 5), "secondary": RGBColor(255, 171, 0)},
        "Monochrome": {"primary": RGBColor(52, 73, 94), "secondary": RGBColor(44, 62, 80)}
    }
    return schemes.get(scheme_name, schemes["Google Blue"])

def get_font_size(size_option):
    sizes = {
        "Small (18pt)": {"title": Pt(28), "content": Pt(18)},
        "Medium (24pt)": {"title": Pt(36), "content": Pt(24)},
        "Large (32pt)": {"title": Pt(44), "content": Pt(32)}
    }
    return sizes.get(size_option, sizes["Medium (24pt)"])

def generate_presentation_content_gemini(topic, description, num_slides, style, audience, include_images, api_key):
    """Generate presentation content using Gemini 2.0 Flash API"""
    
    if not api_key:
        raise ValueError("Please provide your Gemini API key")
    
    # Create the prompt
    prompt = f"""
    Create a comprehensive PowerPoint presentation about "{topic}".
    
    Requirements:
    - Number of slides: {num_slides}
    - Style: {style}
    - Target audience: {audience}
    - Additional details: {description if description else "None"}
    - Include image suggestions: {include_images}
    
    Please provide the response in valid JSON format exactly like this:
    {{
        "title": "Presentation Title",
        "slides": [
            {{
                "slide_number": 1,
                "title": "Slide Title",
                "content": [
                    "Bullet point 1",
                    "Bullet point 2",
                    "Bullet point 3"
                ],
                "image_suggestion": "Description of relevant image",
                "notes": "Speaker notes for this slide"
            }}
        ]
    }}
    
    Rules:
    1. Create exactly {num_slides} slides
    2. Each slide must have 3-5 bullet points
    3. Make content appropriate for {audience} audience
    4. Use {style} presentation style
    5. Include speaker notes for each slide
    6. {"Include relevant image suggestions" if include_images else "Do not include image suggestions"}
    7. Return ONLY valid JSON, no other text
    
    The presentation should be informative, well-structured, and engaging.
    """
    
    # Prepare the request
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={api_key}"
    
    headers = {
        'Content-Type': 'application/json'
    }
    
    data = {
        "contents": [
            {
                "parts": [
                    {
                        "text": prompt
                    }
                ]
            }
        ],
        "generationConfig": {
            "temperature": 0.7,
            "topK": 40,
            "topP": 0.95,
            "maxOutputTokens": 4096,
        }
    }
    
    # Make API call
    try:
        response = requests.post(url, headers=headers, json=data)
        response.raise_for_status()
        
        response_data = response.json()
        
        # Extract the generated text
        if 'candidates' in response_data and len(response_data['candidates']) > 0:
            generated_text = response_data['candidates'][0]['content']['parts'][0]['text']
            
            # Clean and parse JSON
            # Remove markdown code blocks if present
            cleaned_text = re.sub(r'```json\s*', '', generated_text)
            cleaned_text = re.sub(r'```\s*$', '', cleaned_text)
            cleaned_text = cleaned_text.strip()
            
            # Find JSON in the response
            start_idx = cleaned_text.find('{')
            end_idx = cleaned_text.rfind('}') + 1
            
            if start_idx != -1 and end_idx != -1:
                json_str = cleaned_text[start_idx:end_idx]
                presentation_data = json.loads(json_str)
                
                # Validate the structure
                if 'title' in presentation_data and 'slides' in presentation_data:
                    return presentation_data
            
            # If JSON parsing fails, create a fallback structure
            raise json.JSONDecodeError("Could not parse JSON", generated_text, 0)
            
        else:
            raise Exception("No response generated from Gemini API")
            
    except requests.exceptions.RequestException as e:
        raise Exception(f"API request failed: {str(e)}")
    except json.JSONDecodeError as e:
        # Create a fallback presentation structure
        return create_fallback_presentation(topic, num_slides, generated_text if 'generated_text' in locals() else "")
    except Exception as e:
        raise Exception(f"Error generating content: {str(e)}")

def create_fallback_presentation(topic, num_slides, raw_text=""):
    """Create a fallback presentation structure"""
    slides = []
    
    # Create basic slides
    slide_titles = [
        "Introduction",
        "Overview",
        "Key Points",
        "Benefits",
        "Implementation",
        "Challenges",
        "Solutions",
        "Future Outlook",
        "Case Studies",
        "Best Practices",
        "Recommendations",
        "Conclusion",
        "Questions & Discussion"
    ]
    
    for i in range(num_slides):
        slide_num = i + 1
        title = slide_titles[i % len(slide_titles)]
        if i == 0:
            title = f"Introduction to {topic}"
        elif i == num_slides - 1:
            title = "Conclusion"
        
        slides.append({
            "slide_number": slide_num,
            "title": title,
            "content": [
                f"Key point about {topic}",
                f"Important aspect to consider",
                f"Relevant information for {title.lower()}",
                f"Additional insights"
            ],
            "image_suggestion": f"Relevant image for {title.lower()}",
            "notes": f"Speaker notes for slide {slide_num} about {title.lower()}"
        })
    
    return {
        "title": f"{topic} - Presentation",
        "slides": slides
    }

def create_powerpoint(presentation_data, color_scheme, font_sizes, template_style):
    """Create PowerPoint presentation from generated content"""
    
    prs = Presentation()
    colors = get_color_scheme(color_scheme)
    fonts = get_font_size(font_sizes)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = presentation_data["title"]
    subtitle.text = f"Generated by Gemini 2.0 Flash • {datetime.now().strftime('%B %Y')}"
    
    # Style title slide
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = fonts["title"]
        paragraph.font.color.rgb = colors["primary"]
        paragraph.font.bold = True
    
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = fonts["content"]
        paragraph.font.color.rgb = colors["secondary"]
    
    # Content slides
    for slide_data in presentation_data["slides"]:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = slide_data["title"]
        
        # Style title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = fonts["title"]
            paragraph.font.color.rgb = colors["primary"]
            paragraph.font.bold = True
        
        # Content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, point in enumerate(slide_data["content"]):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = fonts["content"]
            p.font.color.rgb = RGBColor(64, 64, 64)
        
        # Add speaker notes if available
        if "notes" in slide_data:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]
    
    return prs

# Generate button
if st.button(" Generate Presentation with Gemini", type="primary"):
    if not topic:
        st.error("Please enter a presentation topic!")
    elif not api_key:
        st.error("Please provide your Gemini API key in the sidebar!")
    else:
        with st.spinner(" Gemini 2.0 Flash is creating your presentation..."):
            try:
                # Generate content
                presentation_data = generate_presentation_content_gemini(
                    topic, description, num_slides, presentation_style, 
                    audience, include_images, api_key
                )
                
                st.session_state.generated_content = presentation_data
                
                # Create PowerPoint
                pptx = create_powerpoint(
                    presentation_data, color_scheme, font_size, template_style
                )
                
                # Save to bytes
                pptx_io = io.BytesIO()
                pptx.save(pptx_io)
                pptx_io.seek(0)
                st.session_state.pptx_file = pptx_io.getvalue()
                
                st.success("✅ Presentation generated successfully with Gemini 2.0 Flash!")
                
            except Exception as e:
                st.error(f"❌ Error generating presentation: {str(e)}")

# Display generated content
if st.session_state.generated_content:
    st.header(" Generated Content Preview")
    
    data = st.session_state.generated_content
    st.subheader(f" {data['title']}")
    
    # Display slides in a more compact format
    for i, slide in enumerate(data['slides'], 1):
        with st.expander(f"Slide {i}: {slide['title']}", expanded=i==1):
            col_content, col_info = st.columns([3, 1])
            
            with col_content:
                st.write("**Content:**")
                for point in slide['content']:
                    st.write(f"• {point}")
            
            with col_info:
                if 'image_suggestion' in slide and slide['image_suggestion']:
                    st.write("** Image Idea:**")
                    st.caption(slide['image_suggestion'])
                
                if 'notes' in slide and slide['notes']:
                    st.write("** Notes:**")
                    st.caption(slide['notes'])

# Download button
if st.session_state.pptx_file:
    st.header(" Download Your Presentation")
    
    filename = f"{topic.replace(' ', '_')}_gemini_presentation.pptx" if topic else "gemini_ai_presentation.pptx"
    
    col_download, col_info = st.columns([1, 2])
    
    with col_download:
        st.download_button(
            label=" Download PowerPoint File",
            data=st.session_state.pptx_file,
            file_name=filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    
    with col_info:
        st.info(" **Your presentation is ready!**\n- Open with PowerPoint, Google Slides, or LibreOffice\n- All speaker notes are included\n- Fully editable and customizable")

# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: gray; padding: 2rem;">
    <p> Powered by <strong>Google Gemini 2.0 Flash</strong> • Built with Streamlit</p>
    <p> Get your free API key at <a href="https://makersuite.google.com/app/apikey" target="_blank">Google AI Studio</a></p>
</div>
""", unsafe_allow_html=True)















